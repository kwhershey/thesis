from pptx import Presentation
from pptx.util import Inches

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
blank_slide_layout = prs.slide_layouts[6]

with open('thesis.log','r') as f:
    for l in f:
        if l.split(' ')[0]=='Chapter':
            slide = prs.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            title.text = l
        elif l.split(' ')[0]=='<use':
            slide = prs.slides.add_slide(blank_slide_layout)
            pic = slide.shapes.add_picture(l.split(' ')[1][:-2], Inches(1), Inches(1))
prs.save('images/figs.pptx')
